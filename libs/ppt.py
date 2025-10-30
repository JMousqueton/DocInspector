import os
from pptx import Presentation

def is_pptx_file(filename):
    # True if file is .pptx and starts with PK (zip signature)
    return filename.lower().endswith('.pptx') and _has_zip_sig(filename)

def _has_zip_sig(filename):
    try:
        with open(filename, "rb") as f:
            return f.read(2) == b"PK"
    except Exception:
        return False

def extract_pptx_metadata(pptx_file):
    prs = Presentation(pptx_file)
    core = prs.core_properties
    meta = [
        ["title", core.title or ""],
        ["subject", core.subject or ""],
        ["creator", core.author or ""],
        ["keywords", core.keywords or ""],
        ["description", core.comments or ""],
        ["last_modified_by", core.last_modified_by or ""],
        ["revision", core.revision or ""],
        ["created", str(core.created) if core.created else ""],
        ["modified", str(core.modified) if core.modified else ""],
        ["category", core.category or ""],
        ["content_status", core.content_status or ""],
        ["identifier", core.identifier or ""],
        ["language", core.language or ""],
        ["version", core.version or ""]
    ]
    # Custom properties (if any)
    if hasattr(prs, 'custom_properties'):
        for key in prs.custom_properties:
            meta.append([f"custom_{key}", str(prs.custom_properties[key])])
    # Optionally extract the template property (from app.xml)
    template = get_pptx_template_name(pptx_file)
    if template:
        meta.append(["template", template])
    # Optionally extract theme names
    theme_names = get_pptx_theme_names(pptx_file)
    if theme_names:
        meta.append(["themes", ", ".join(theme_names)])
    return meta

def get_pptx_basic_info(pptx_file):
    from libs.shared import human_readable_size
    file_size = os.path.getsize(pptx_file)
    prs = Presentation(pptx_file)
    meta = extract_pptx_metadata(pptx_file)
    num_slides = len(prs.slides)
    images = extract_pptx_images(pptx_file)
    links = extract_pptx_links(prs)
    comments = extract_pptx_comments(pptx_file)

    slides_with_notes = []
    notes_texts = {}
    for i, slide in enumerate(prs.slides):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            slides_with_notes.append(i + 1)  # 1-based
            notes_texts[i + 1] = slide.notes_slide.notes_text_frame.text.strip()
    num_slides_with_notes = len(slides_with_notes)

    # Custom XML
    custom_xml_parts = extract_custom_xml_parts(pptx_file)
    has_macros = has_vba_macros(pptx_file)

    return {
        "file_size_bytes": file_size,
        "file_size_human": human_readable_size(file_size),
        "num_slides": num_slides,
        "num_slides_with_notes": num_slides_with_notes,
        "slides_with_notes": slides_with_notes,
        "notes_texts": notes_texts,
        "meta": meta,
        "images": images,
        "links": links,
        "comments": comments,
        "custom_xml_parts": custom_xml_parts,
        "has_vba_macros": has_macros,
    }


def extract_pptx_images(pptx_file):
    # Extract image names from the pptx zip (ppt/media/*)
    from zipfile import ZipFile
    img_names = []
    with ZipFile(pptx_file) as pptx_zip:
        img_names = [name for name in pptx_zip.namelist() if name.startswith('ppt/media/')]
    return img_names

def extract_pptx_links(prs):
    # Extract hyperlinks from shapes, skip group shapes
    links = set()
    from pptx.shapes.group import GroupShape
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.hyperlink.address:
                            links.add(run.hyperlink.address)
            # Only try click_action on non-group shapes
            if not isinstance(shape, GroupShape):
                try:
                    if hasattr(shape, "click_action") and shape.click_action.hyperlink.address:
                        links.add(shape.click_action.hyperlink.address)
                except Exception:
                    pass
    return sorted(links)

def extract_pptx_comments(pptx_file):
    # Parse comments from ppt/comments*.xml
    from zipfile import ZipFile
    import xml.etree.ElementTree as ET
    comments = []
    with ZipFile(pptx_file) as pptx_zip:
        for name in pptx_zip.namelist():
            if name.startswith('ppt/comments') and name.endswith('.xml'):
                with pptx_zip.open(name) as f:
                    xml = f.read()
                    root = ET.fromstring(xml)
                    for comment in root.iter('{http://schemas.openxmlformats.org/presentationml/2006/main}cm'):
                        author = comment.attrib.get('authorId', '')
                        date = comment.attrib.get('dt', '')
                        text = ''.join(child.text or '' for child in comment.iter('{http://schemas.openxmlformats.org/presentationml/2006/main}t'))
                        comments.append({'author': author, 'date': date, 'text': text})
    return comments

def get_pptx_template_name(pptx_file):
    # Try to extract template property from docProps/app.xml
    from zipfile import ZipFile
    import xml.etree.ElementTree as ET
    try:
        with ZipFile(pptx_file) as pptx_zip:
            with pptx_zip.open("docProps/app.xml") as f:
                tree = ET.parse(f)
                root = tree.getroot()
                ns = {'ap': 'http://schemas.openxmlformats.org/officeDocument/2006/extended-properties'}
                template_elem = root.find('ap:Template', ns)
                if template_elem is not None:
                    return template_elem.text
    except Exception:
        pass
    return ""

def get_pptx_theme_names(pptx_file):
    # Returns a list of theme names (may be generic, e.g., 'Office Theme')
    from pptx import Presentation
    prs = Presentation(pptx_file)
    names = []
    for slide_master in prs.slide_masters:
        try:
            names.append(slide_master.theme.name)
        except Exception:
            pass
    return names


def extract_custom_xml_parts(pptx_file):
    """
    Returns a list of custom XML part filenames and (optionally) their contents.
    """
    from zipfile import ZipFile
    xml_parts = []
    with ZipFile(pptx_file) as pptx_zip:
        for name in pptx_zip.namelist():
            if name.startswith("customXml/") and name.endswith(".xml"):
                # You can just list the name, or extract the content:
                try:
                    with pptx_zip.open(name) as f:
                        content = f.read().decode("utf-8", errors="replace")
                    xml_parts.append({"filename": name, "content": content})
                except Exception:
                    xml_parts.append({"filename": name, "content": "(unreadable)"})
    return xml_parts

def has_vba_macros(pptx_file):
    """
    Returns True if ppt/vbaProject.bin is present, indicating VBA macros.
    """
    from zipfile import ZipFile
    with ZipFile(pptx_file) as pptx_zip:
        return "ppt/vbaProject.bin" in pptx_zip.namelist()
